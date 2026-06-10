
from fastapi import FastAPI, UploadFile, File, HTTPException, Request
from fastapi.middleware.cors import CORSMiddleware
from dotenv import load_dotenv
import os
import uuid
import shutil
import pytesseract
from pdf2image import convert_from_path
from PIL import Image
import re
from transformers import pipeline

# --------- LOAD .ENV ------------
load_dotenv()

# --------- OCR CONFIG -----------
pytesseract.pytesseract.tesseract_cmd = r"C:\Users\acer\AppData\Local\Programs\Tesseract-OCR\tesseract.exe"
POPPLER_PATH = r"C:\Users\acer\Downloads\Release-25.11.0-0\poppler-25.11.0\Library\bin"

# --------- UPLOAD DIRECTORY -----
UPLOAD_DIR = os.path.join("data", "uploads")
os.makedirs(UPLOAD_DIR, exist_ok=True)

# --------- FASTAPI APP ----------
app = FastAPI(title="AI Document Intelligence System", version="7.2")

# --------- CORS FIX --------------
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.options("/{rest_of_path:path}")
async def preflight(rest_of_path: str):
    return {}

# --------- SAVE FILE ------------
async def save_file(file: UploadFile):
    file_id = str(uuid.uuid4())[:8]
    filename = f"{file_id}_{file.filename}"
    path = os.path.join(UPLOAD_DIR, filename)

    with open(path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    return path, file_id, filename

# --------- EXTRACTORS -----------
def extract_text_from_image(path):
    img = Image.open(path).convert("L")
    text = pytesseract.image_to_string(img)

    print("\n========== OCR OUTPUT ==========")
    print(text)
    print("================================\n")

    return text


def extract_text_from_pil_image(pil_img):
    img = pil_img.convert("L")
    text = pytesseract.image_to_string(img)

    print("\n========== OCR OUTPUT ==========")
    print(text)
    print("================================\n")

    return text
def extract_image(path):
    return extract_text_from_image(path)

def extract_pdf(path):
    text = ""

    try:
        pages = convert_from_path(path, dpi=300, poppler_path=POPPLER_PATH)

        for p in pages:
            text += extract_text_from_pil_image(p) + "\n"

        return text

    except Exception as e:
        print("PDF OCR failed:", e)

    try:
        import PyPDF2

        with open(path, "rb") as f:
            reader = PyPDF2.PdfReader(f)

            for page in reader.pages:
                text += page.extract_text() or ""

        return text

    except Exception as e:
        print("PyPDF2 extraction failed:", e)
        return ""

def extract_docx(path):
    from docx import Document
    doc = Document(path)
    return "\n".join([p.text for p in doc.paragraphs])

def extract_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_excel(path):
    import pandas as pd
    xl = pd.ExcelFile(path)
    text = ""
    for sheet in xl.sheet_names:
        df = pd.read_excel(path, sheet)
        text += f"--- Sheet: {sheet} ---\n"
        text += df.to_string() + "\n"
    return text

def extract_csv(path):
    import pandas as pd
    df = pd.read_csv(path)
    return df.to_string()

def extract_txt(path):
    try:
        return open(path, "r", encoding="utf-8").read()
    except:
        return open(path, "r", encoding="latin-1").read()

# --------- CLEAN TEXT ----------
def clean_text(text):
    if not text:
        return ""

    # weird unicode chars remove
    text = re.sub(r"[^\w\s.,:/&()\-']", " ", text)

    # multiple spaces
    text = re.sub(r"\s+", " ", text)

    # multiple newlines
    text = re.sub(r"\n+", "\n", text)

    return text.strip()

# --------- CLASSIFIER ----------
def classify(text):
    t = text.lower()
    if "aadhaar" in t or re.search(r"\b\d{4}\s\d{4}\s\d{4}\b", t):
        return "aadhaar_card"
    if re.search(r"[A-Z]{5}[0-9]{4}[A-Z]", t):
        return "pan_card"
    if "resume" in t:
        return "resume"
    if "invoice" in t:
        return "invoice"
    if "certificate" in t:
        return "certificate"
    return "unknown"

# --------- FIELD EXTRACTION -----
def extract_fields(text):
    return {
        "emails": list(set(re.findall(r"[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[A-Z|a-z]{2,}", text))),
        "phones": list(set(re.findall(r"\b\d{10}\b", text))),
        "aadhaar_numbers": list(set(re.findall(r"\b\d{4}\s\d{4}\s\d{4}\b", text))),
        "pan_numbers": list(set(re.findall(r"[A-Z]{5}[0-9]{4}[A-Z]", text))),
        "dates": list(set(re.findall(r"\b\d{1,2}[/-]\d{1,2}[/-]\d{2,4}\b", text))),
        "amounts": list(set(re.findall(r"₹\s?\d+(?:,\d{3})*", text))),
    }

def filter_fields(fields):
    return {k: v for k, v in fields.items() if v}

# --------- YOUR ORIGINAL SUMMARIZER ---------
def ai_summary(text, doc_type):
    text = text.strip()

    if len(text) < 30:
        return "Not enough content to generate meaningful summary."

    clean = re.sub(r"[¢@]+", "", text)
    clean = re.sub(r"\s+", " ", clean).strip()

    try:
        ai = summarizer(
            clean[:2500],
            max_length=120,
            min_length=40,
            do_sample=False
        )[0]["summary_text"]
    except:
        ai = clean[:700]

    return f"""
Document Type: {doc_type}

Summary:
{ai}
""".strip()

# --------- SAFE EXTRA FEATURES -----
def confidence_score(doc_type, fields):
    score = 50
    if doc_type != "unknown":
        score += 25
    if len(fields.keys()) >= 2:
        score += 25
    return f"{score}%"

def smart_doc_name(doc_type):
    return f"{doc_type}_document.pdf"

def fraud_check(text):
    return any(x in text.lower() for x in ["fake", "forged", "tampered"])

def redact_sensitive(text):
    text = re.sub(r"\b\d{4}\s\d{4}\s\d{4}\b", "XXXX XXXX XXXX", text)
    text = re.sub(r"[A-Z]{5}[0-9]{4}[A-Z]", "XXXXX0000X", text)
    return text

# --------- MAIN ENDPOINT ----------
@app.post("/upload")
async def upload(file: UploadFile = File(...), request: Request = None):
    user_id = request.headers.get("x-user-id", "guest")

    ext = file.filename.lower().split(".")[-1]
    path, file_id, filename = await save_file(file)

    if ext in ["png", "jpg", "jpeg"]:
        extracted = extract_image(path)
    elif ext == "pdf":
        extracted = extract_pdf(path)
    elif ext == "docx":
        extracted = extract_docx(path)
    elif ext == "pptx":
        extracted = extract_pptx(path)
    elif ext in ["xlsx", "xls"]:
        extracted = extract_excel(path)
    elif ext == "csv":
        extracted = extract_csv(path)
    elif ext == "txt":
        extracted = extract_txt(path)
    else:
        raise HTTPException(400, "Unsupported file type")

    cleaned = clean_text(extracted)
    doc_type = classify(cleaned)
    fields = filter_fields(extract_fields(cleaned))
    summary = ai_summary(cleaned, doc_type)

    return {
        "user_id": user_id,
        "file_id": file_id,
        "smart_name": smart_doc_name(doc_type),
        "document_type": doc_type,
        "confidence_score": confidence_score(doc_type, fields),
        "word_count": len(cleaned.split()),
        "fraud_flag": fraud_check(cleaned),
        "summary": summary,
        "fields": fields,
        "redacted_text": redact_sensitive(cleaned),
        "message": "Document processed successfully ✅"
    }

@app.get("/")
def home():
    return {"message": "AI Document Intelligence System running ✅"}
